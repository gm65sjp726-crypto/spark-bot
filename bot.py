import sqlite3
from datetime import datetime
from telegram import Update
from telegram.ext import ApplicationBuilder, CommandHandler, MessageHandler, ContextTypes, filters
from pptx import Presentation

# ====== НАСТРОЙКИ ======
TOKEN = "7719110832:AAHkFBu0QbIly_gQgLDtbniO8XAxQ6t0GZE"

MONTHLY_SPARKS = 400
PRESENTATION_COST = 80
FREE_SLIDE_LIMIT = 7

# ====== БАЗА ======
conn = sqlite3.connect("bot.db", check_same_thread=False)
cursor = conn.cursor()

cursor.execute("""
CREATE TABLE IF NOT EXISTS users (
    user_id INTEGER PRIMARY KEY,
    sparks INTEGER,
    last_reset TEXT
)
""")

cursor.execute("""
CREATE TABLE IF NOT EXISTS states (
    user_id INTEGER PRIMARY KEY,
    topic TEXT,
    style TEXT,
    font TEXT,
    slides INTEGER
)
""")

conn.commit()

# ====== SPARKS ======
def get_user(user_id):
    cursor.execute("SELECT * FROM users WHERE user_id=?", (user_id,))
    user = cursor.fetchone()
    if not user:
        cursor.execute(
            "INSERT INTO users VALUES (?, ?, ?)",
            (user_id, MONTHLY_SPARKS, datetime.now().isoformat())
        )
        conn.commit()
        return get_user(user_id)
    return user

def has_enough_sparks(user_id):
    get_user(user_id)
    cursor.execute("SELECT sparks FROM users WHERE user_id=?", (user_id,))
    return cursor.fetchone()[0] >= PRESENTATION_COST

def spend_sparks(user_id):
    cursor.execute(
        "UPDATE users SET sparks = sparks - ? WHERE user_id=?",
        (PRESENTATION_COST, user_id)
    )
    conn.commit()

# ====== PPT ======
def create_presentation(topic, style, font, slides_count, filename):
    prs = Presentation()
    titles = ["Введение", "Проблема", "Решение", "Преимущества", "Вывод"]
    slides_count = min(slides_count, FREE_SLIDE_LIMIT)

    for i in range(slides_count):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = titles[i % len(titles)]
        slide.placeholders[1].text = f"{titles[i % len(titles)]} по теме: {topic}"

        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    p.font.name = font

    prs.save(filename)

# ====== BOT ======
async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    get_user(update.message.from_user.id)
    await update.message.reply_text("⚡ Напиши тему презентации")

async def handle_topic(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.message.from_user.id
    topic = update.message.text

    cursor.execute(
        "INSERT OR REPLACE INTO states VALUES (?, ?, ?, ?, ?)",
        (user_id, topic, "бизнес", "Arial", 5)
    )
    conn.commit()

    await update.message.reply_text(
        "Напиши параметры или 'нет'\n"
        "Пример:\nСтиль: бизнес\nШрифт: Arial\nСлайдов: 5"
    )

async def handle_params(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.message.from_user.id
    text = update.message.text.lower()

    cursor.execute("SELECT * FROM states WHERE user_id=?", (user_id,))
    state = cursor.fetchone()
    if not state:
        return

    style, font, slides = state[2], state[3], state[4]

    if "шрифт" in text:
        font = text.split(":")[-1].strip()
    if "слайдов" in text:
        try:
            slides = int(text.split(":")[-1].strip())
        except:
            pass

    if not has_enough_sparks(user_id):
        await update.message.reply_text("⛔ Sparks закончились")
        return

    spend_sparks(user_id)
    filename = f"{user_id}.pptx"
    create_presentation(state[1], style, font, slides, filename)

    await update.message.reply_document(open(filename, "rb"))

app = ApplicationBuilder().token(TOKEN).build()
app.add_handler(CommandHandler("start", start))
app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_topic))
app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_params))
app.run_polling()